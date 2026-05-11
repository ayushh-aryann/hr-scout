"""Generate HR Scout PowerPoint presentation using python-pptx."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand colours ──────────────────────────────────────────────────────────────
INK   = RGBColor(0x0a, 0x0a, 0x0b)
YELLOW= RGBColor(0xff, 0xd8, 0x4d)
WHITE = RGBColor(0xff, 0xff, 0xff)
GREY  = RGBColor(0x6b, 0x72, 0x80)
GREEN = RGBColor(0x16, 0xa3, 0x4a)
AMBER = RGBColor(0xd9, 0x77, 0x06)
RED   = RGBColor(0xdc, 0x26, 0x26)
LIGHT = RGBColor(0xf9, 0xfa, 0xfb)
DARK2 = RGBColor(0x1a, 0x1a, 0x1d)

W = Inches(13.33)
H = Inches(7.5)

SCREENSHOTS = Path("assets/screenshots")
OUT = Path("presentation/HR_Scout_Presentation.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p


def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.15), DARK2)
    add_textbox(slide, title,
                Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
                font_size=28, bold=True, color=YELLOW)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.7), Inches(12), Inches(0.38),
                    font_size=13, color=GREY)


def add_screenshot(slide, name, left, top, width, height=None):
    path = SCREENSHOTS / f"{name}.png"
    if not path.exists():
        return
    if height:
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(path), left, top, width)


def add_accent_line(slide, top_offset=Inches(1.15)):
    add_rect(slide, 0, top_offset, W, Pt(3), YELLOW)


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, INK)

    # Yellow accent strip left
    add_rect(slide, 0, 0, Inches(0.35), H, YELLOW)

    # Large title
    add_textbox(slide, "HR Scout",
                Inches(1.0), Inches(1.6), Inches(11), Inches(1.5),
                font_size=72, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)

    add_textbox(slide, "AI-Powered Candidate Shortlisting Agent",
                Inches(1.0), Inches(3.0), Inches(11), Inches(0.7),
                font_size=26, bold=False, color=WHITE)

    add_textbox(slide, "Automated JD Parsing  •  Resume Ingestion  •  Semantic Scoring  •  Human-in-the-Loop",
                Inches(1.0), Inches(3.75), Inches(11.5), Inches(0.5),
                font_size=14, color=GREY)

    add_textbox(slide, "Built by Ayush Aryan",
                Inches(1.0), Inches(6.6), Inches(5), Inches(0.5),
                font_size=13, color=GREY)

    add_textbox(slide, "2026",
                Inches(11.5), Inches(6.6), Inches(1.5), Inches(0.5),
                font_size=13, color=GREY, align=PP_ALIGN.RIGHT)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "The Problem", "Manual resume screening is broken")
    add_accent_line(slide)

    items = [
        ("250+", "resumes per job posting on average"),
        ("23 hrs", "average recruiter time per hire on screening alone"),
        ("37%",  "of qualified candidates rejected due to keyword-only ATS"),
        ("High bias", "inconsistent scoring across HR reviewers"),
        ("No audit",  "manual processes leave no traceable decision trail"),
    ]

    for i, (stat, desc) in enumerate(items):
        top = Inches(1.5) + i * Inches(1.0)
        add_rect(slide, Inches(0.5), top, Inches(2.2), Inches(0.72),
                 DARK2, YELLOW)
        add_textbox(slide, stat,
                    Inches(0.55), top + Pt(4), Inches(2.1), Inches(0.45),
                    font_size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    Inches(2.9), top + Pt(16), Inches(9.8), Inches(0.4),
                    font_size=15, color=WHITE)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "The Solution", "HR Scout — end-to-end AI shortlisting pipeline")
    add_accent_line(slide)

    boxes = [
        ("JD Parser",       "Extract structured requirements from any job description"),
        ("Resume Ingestion", "PDF, DOCX, and LinkedIn JSON in one drag-drop"),
        ("Semantic Scoring", "5-dimension weighted rubric with embeddings"),
        ("Shortlist Reports","JSON, HTML, and PDF — ready for sharing"),
        ("Human Override",   "Audit-trailed corrections with full justification log"),
    ]

    cols = [(Inches(0.4), Inches(2.6)), (Inches(4.65), Inches(2.6)), (Inches(8.9), Inches(2.6))]
    rows = [Inches(1.4), Inches(4.3)]

    positions = [
        (Inches(0.4),  Inches(1.4)),
        (Inches(4.65), Inches(1.4)),
        (Inches(8.9),  Inches(1.4)),
        (Inches(1.5),  Inches(4.25)),
        (Inches(6.5),  Inches(4.25)),
    ]
    widths = [Inches(3.8), Inches(3.8), Inches(3.8), Inches(4.5), Inches(4.5)]

    for (lbl, desc), (left, top), w in zip(boxes, positions, widths):
        add_rect(slide, left, top, w, Inches(2.35), DARK2, YELLOW)
        add_textbox(slide, lbl, left + Inches(0.15), top + Inches(0.15),
                    w - Inches(0.3), Inches(0.5), font_size=16, bold=True, color=YELLOW)
        add_textbox(slide, desc, left + Inches(0.15), top + Inches(0.65),
                    w - Inches(0.3), Inches(1.4), font_size=13, color=WHITE)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "System Architecture", "Custom pipeline — no LangChain overhead")
    add_accent_line(slide)

    # Pipeline flow boxes
    stages = [
        ("Input Layer",    "PDF/DOCX/JSON\nLinkedIn Export\nJD Text"),
        ("Parse Layer",    "JD Parser\nResume Parser\nLinkedIn Parser"),
        ("Score Layer",    "Embeddings\nRubric Engine\nLLM Scoring"),
        ("Output Layer",   "JSON Report\nHTML Report\nPDF Report"),
    ]

    for i, (title, detail) in enumerate(stages):
        left = Inches(0.3 + i * 3.2)
        top = Inches(1.6)
        w = Inches(2.9)
        h = Inches(4.6)
        color = YELLOW if i == 2 else DARK2
        text_color = INK if i == 2 else WHITE

        add_rect(slide, left, top, w, h, color)
        add_textbox(slide, title, left + Inches(0.1), top + Inches(0.1),
                    w - Inches(0.2), Inches(0.5), font_size=16, bold=True,
                    color=INK if i == 2 else YELLOW)
        add_textbox(slide, detail, left + Inches(0.15), top + Inches(0.75),
                    w - Inches(0.3), Inches(3.5), font_size=13, color=text_color)

        # Arrow between boxes
        if i < 3:
            add_textbox(slide, "→",
                        left + w, top + Inches(2.0), Inches(0.3), Inches(0.5),
                        font_size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # FastAPI + React note
    add_textbox(slide,
                "FastAPI REST backend  •  React 18 SPA (CDN, no build step)  •  Pydantic validation  •  sentence-transformers embeddings",
                Inches(0.3), Inches(6.8), Inches(12.5), Inches(0.45),
                font_size=11, color=GREY, align=PP_ALIGN.CENTER)


def slide_ai_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "AI Workflow", "LLM + semantic embeddings + heuristic fallback")
    add_accent_line(slide)

    steps = [
        ("1", "Sanitize",   "Strip prompt-injection patterns\n12 regex rules"),
        ("2", "Parse JD",   "LLM → structured JDProfile\nHeuristic fallback"),
        ("3", "Parse CVs",  "PDF/DOCX text extraction\nLLM → CandidateProfile"),
        ("4", "Embed",      "sentence-transformers\ncosine similarity"),
        ("5", "Score",      "5-dimension rubric\nWeighted 30/25/15/20/10"),
        ("6", "Rank",       "Sort by total score\nRecommendation threshold"),
    ]

    for i, (num, title, detail) in enumerate(steps):
        col = i % 3
        row = i // 3
        left = Inches(0.4 + col * 4.3)
        top = Inches(1.5 + row * 2.7)
        w = Inches(3.9)
        h = Inches(2.4)

        add_rect(slide, left, top, w, h, DARK2, YELLOW)
        # Number badge
        add_rect(slide, left + Inches(0.15), top + Inches(0.15),
                 Inches(0.45), Inches(0.45), YELLOW)
        add_textbox(slide, num, left + Inches(0.15), top + Inches(0.1),
                    Inches(0.45), Inches(0.45), font_size=16, bold=True,
                    color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, left + Inches(0.75), top + Inches(0.15),
                    w - Inches(0.9), Inches(0.5), font_size=15, bold=True, color=YELLOW)
        add_textbox(slide, detail, left + Inches(0.2), top + Inches(0.75),
                    w - Inches(0.4), Inches(1.4), font_size=12, color=WHITE)

    add_textbox(slide,
                "No API key required — heuristic mode enables full demo without LLM",
                Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.38),
                font_size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def slide_scoring(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "Scoring Rubric", "5 dimensions — tamper-proof via Pydantic validators")
    add_accent_line(slide)

    rows = [
        ("Skills Match",          "30%", 30, "Semantic overlap + exact keyword match"),
        ("Experience Relevance",  "25%", 25, "Years, seniority, domain alignment"),
        ("Education & Certs",     "15%", 15, "Degree level, relevant certifications"),
        ("Project Portfolio",     "20%", 20, "Relevant projects, GitHub, impact"),
        ("Communication Quality", "10%", 10, "Writing clarity, formatting, structure"),
    ]

    # Header row
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.4), Inches(0.45), DARK2)
    for lbl, left_offset in [("Dimension", 0.4), ("Weight", 5.2), ("Bar", 6.3), ("What it measures", 8.6)]:
        add_textbox(slide, lbl, Inches(left_offset), Inches(1.33),
                    Inches(2.5), Inches(0.35), font_size=11, bold=True,
                    color=GREY)

    for i, (name, weight, pct, desc) in enumerate(rows):
        top = Inches(1.8) + i * Inches(0.95)
        bg = DARK2 if i % 2 == 0 else RGBColor(0x12, 0x12, 0x14)
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(0.85), bg)

        add_textbox(slide, name, Inches(0.55), top + Inches(0.2),
                    Inches(4.5), Inches(0.5), font_size=15, bold=True, color=WHITE)
        add_textbox(slide, weight, Inches(5.2), top + Inches(0.2),
                    Inches(0.9), Inches(0.5), font_size=20, bold=True,
                    color=YELLOW, align=PP_ALIGN.CENTER)

        # Progress bar
        bar_w = Inches(1.8 * pct / 30)
        add_rect(slide, Inches(6.3), top + Inches(0.3), Inches(1.8), Inches(0.25),
                 RGBColor(0x2d, 0x2d, 0x2f))
        add_rect(slide, Inches(6.3), top + Inches(0.3), bar_w, Inches(0.25), YELLOW)

        add_textbox(slide, desc, Inches(8.6), top + Inches(0.2),
                    Inches(4.0), Inches(0.5), font_size=12, color=GREY)

    add_textbox(slide,
                "Pydantic model_validator recomputes weighted_contribution after every LLM response — scores cannot be inflated",
                Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.38),
                font_size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def slide_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "Security Mitigations", "Defense-in-depth for production use")
    add_accent_line(slide)

    items = [
        ("Prompt Injection Defense",
         "12 regex patterns block 'ignore previous instructions', jailbreaks, system tags"),
        ("PII Masking",
         "Email, phone, SSN, passport masked in all audit logs before persistence"),
        ("API Key Auth",
         "X-API-Key header; .env-only storage; never hardcoded in source"),
        ("Input Validation",
         "Pydantic strict types on every model field; scores clamped to 0–10"),
        ("Audit Trail",
         "Append-only JSONL log; all overrides signed with HR user + timestamp"),
    ]

    for i, (title, desc) in enumerate(items):
        top = Inches(1.45 + i * 1.1)
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(0.9), DARK2)
        # Icon placeholder
        add_rect(slide, Inches(0.45), top + Inches(0.15),
                 Inches(0.6), Inches(0.6), YELLOW)
        add_textbox(slide, str(i + 1), Inches(0.45), top + Inches(0.12),
                    Inches(0.6), Inches(0.6), font_size=18, bold=True,
                    color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, Inches(1.2), top + Inches(0.1),
                    Inches(3.8), Inches(0.4), font_size=14, bold=True, color=YELLOW)
        add_textbox(slide, desc, Inches(1.2), top + Inches(0.48),
                    Inches(11.0), Inches(0.38), font_size=12, color=GREY)


def slide_demo_screenshots(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "Demo Screenshots", "Live app — React SPA + FastAPI")
    add_accent_line(slide)

    shots = [
        ("01_dashboard", Inches(0.3),  Inches(1.3), Inches(6.3)),
        ("02_upload_page", Inches(6.8),  Inches(1.3), Inches(6.3)),
        ("04_ranking_table", Inches(0.3),  Inches(4.4), Inches(6.3)),
        ("07_api_docs", Inches(6.8),  Inches(4.4), Inches(6.3)),
    ]

    labels = [
        ("Dashboard / Landing", Inches(0.3),  Inches(2.9)),
        ("Upload & JD Entry",   Inches(6.8),  Inches(2.9)),
        ("Ranking Table",       Inches(0.3),  Inches(6.0)),
        ("API Documentation",   Inches(6.8),  Inches(6.0)),
    ]

    for (name, left, top, width), (lbl, ll, lt) in zip(shots, labels):
        add_screenshot(slide, name, left, top, width, Inches(1.5))
        add_textbox(slide, lbl, ll, lt + Inches(1.52), Inches(6),
                    Inches(0.35), font_size=11, color=GREY, align=PP_ALIGN.CENTER)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "Sample Results", "5 candidates scored against Senior Full-Stack Engineer JD")
    add_accent_line(slide)

    candidates = [
        ("1", "Alice Chen",    "82.8", "HIRE",       GREEN),
        ("2", "Carol Smith",   "71.5", "HIRE",       GREEN),
        ("3", "Bob Kumar",     "65.2", "HIRE",       GREEN),
        ("4", "David Lee",     "51.7", "BORDERLINE", AMBER),
        ("5", "Emma Jones",    "32.4", "NO HIRE",    RED),
    ]

    # Header
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(12.4), Inches(0.45), DARK2)
    for lbl, lft in [("#", 0.55), ("Candidate", 1.1), ("Total Score", 5.8), ("Decision", 9.5)]:
        add_textbox(slide, lbl, Inches(lft), Inches(1.38), Inches(3), Inches(0.35),
                    font_size=11, bold=True, color=GREY)

    for i, (rank, name, score, rec, color) in enumerate(candidates):
        top = Inches(1.85 + i * 1.0)
        bg = DARK2 if i % 2 == 0 else RGBColor(0x12, 0x12, 0x14)
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(0.85), bg)

        add_textbox(slide, rank, Inches(0.55), top + Inches(0.18),
                    Inches(0.5), Inches(0.5), font_size=20, bold=True,
                    color=GREY, align=PP_ALIGN.CENTER)
        add_textbox(slide, name, Inches(1.1), top + Inches(0.18),
                    Inches(4.5), Inches(0.5), font_size=16, bold=True, color=WHITE)
        add_textbox(slide, score + " / 100", Inches(5.8), top + Inches(0.18),
                    Inches(3.5), Inches(0.5), font_size=20, bold=True, color=color)

        # Badge
        add_rect(slide, Inches(9.5), top + Inches(0.2), Inches(2.8), Inches(0.44),
                 RGBColor(max(0, color[0] - 0x99), max(0, color[1] - 0x99), max(0, color[2] - 0x99)))
        add_textbox(slide, rec, Inches(9.5), top + Inches(0.18),
                    Inches(2.8), Inches(0.44), font_size=12, bold=True,
                    color=color, align=PP_ALIGN.CENTER)


def slide_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)
    add_header_bar(slide, "Future Improvements", "Roadmap for production deployment")
    add_accent_line(slide)

    items = [
        ("LinkedIn Scraper Integration",  "Directly pull profiles via LinkedIn API or scraper"),
        ("Vector DB (Pinecone / Weaviate)", "Store all candidate embeddings for cross-role search"),
        ("Bias Audit Module",             "Flag demographic proxies; ensure EEOC compliance"),
        ("Calendar Integration",          "Auto-schedule interviews for HIRE-tier candidates"),
        ("Multi-Tenant SaaS",             "Per-org scoring weights; SSO; usage metering"),
        ("Candidate Portal",              "Self-service profile updates + status tracker"),
    ]

    for i, (title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.4 + col * 6.5)
        top = Inches(1.4 + row * 1.85)
        w = Inches(6.1)
        h = Inches(1.6)

        add_rect(slide, left, top, w, h, DARK2, YELLOW)
        add_textbox(slide, title, left + Inches(0.2), top + Inches(0.12),
                    w - Inches(0.4), Inches(0.5), font_size=14, bold=True, color=YELLOW)
        add_textbox(slide, desc, left + Inches(0.2), top + Inches(0.65),
                    w - Inches(0.4), Inches(0.75), font_size=12, color=WHITE)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK)

    add_rect(slide, 0, 0, Inches(0.35), H, YELLOW)

    add_textbox(slide, "HR Scout",
                Inches(1.0), Inches(1.8), Inches(11), Inches(1.3),
                font_size=64, bold=True, color=YELLOW)

    add_textbox(slide, "Faster hiring. Fairer scoring. Full audit trail.",
                Inches(1.0), Inches(3.0), Inches(11), Inches(0.7),
                font_size=24, color=WHITE)

    add_textbox(slide,
                "FastAPI  •  React 18  •  Claude Sonnet  •  sentence-transformers  •  ReportLab  •  Playwright",
                Inches(1.0), Inches(3.85), Inches(11.5), Inches(0.5),
                font_size=13, color=GREY)

    add_textbox(slide, "github.com/ayusharyan/hr-scout",
                Inches(1.0), Inches(5.5), Inches(7), Inches(0.5),
                font_size=16, color=YELLOW)

    add_textbox(slide, "Built by Ayush Aryan — 2026",
                Inches(1.0), Inches(6.6), Inches(5), Inches(0.5),
                font_size=13, color=GREY)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides...")
    slide_title(prs)
    print("  1/10 Title")
    slide_problem(prs)
    print("  2/10 Problem Statement")
    slide_solution(prs)
    print("  3/10 Solution Overview")
    slide_architecture(prs)
    print("  4/10 Architecture")
    slide_ai_workflow(prs)
    print("  5/10 AI Workflow")
    slide_scoring(prs)
    print("  6/10 Scoring Rubric")
    slide_security(prs)
    print("  7/10 Security Mitigations")
    slide_demo_screenshots(prs)
    print("  8/10 Demo Screenshots")
    slide_results(prs)
    print("  9/10 Sample Results")
    slide_future(prs)
    print(" 10/10 Future Improvements")
    slide_closing(prs)
    print(" 11/11 Closing Slide")

    prs.save(str(OUT))
    print(f"\nPresentation saved to: {OUT}")


if __name__ == "__main__":
    main()
